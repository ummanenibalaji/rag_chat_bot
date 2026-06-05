import os
from datetime import datetime

from langchain_community.document_loaders import (
    PyPDFLoader
)

from langchain_core.documents import (
    Document
)

from ocr_utils import (
    extract_text_from_image
)

from vision_utils import (
    detect_document_type,
    analyze_image
)

from docx import Document as DocxDocument

from pptx import Presentation

import pandas as pd

# =====================================================
# COMMON METADATA
# =====================================================

def build_metadata(
    file_path,
    file_type,
    page=None
):

    filename = os.path.basename(
        file_path
    )

    metadata = {
        "source": filename,
        "file_name": filename,
        "file_type": file_type,
        "document_id": filename.replace(
            ".", "_"
        ),
        "upload_time": str(
            datetime.now()
        )
    }

    if page is not None:

        metadata["page"] = page

    return metadata


# =====================================================
# SCANNED PDF DETECTION
# =====================================================

def is_scanned_pdf(documents: list) -> bool:
    """Return True if PDF has negligible extractable text (scanned/image-based)."""
    if not documents:
        return True
    total_chars = sum(len(doc.page_content.strip()) for doc in documents)
    avg_chars_per_page = total_chars / len(documents)
    return avg_chars_per_page < 120


# =====================================================
# LOAD PDF
# =====================================================

def load_pdf(file_path):
    from ocr_utils import ocr_scanned_pdf

    loader = PyPDFLoader(file_path)
    documents = loader.load()

    filename = os.path.basename(file_path).lower()
    is_drawing = any(
        kw in filename for kw in ["drawing", "cad", "plan", "blueprint"]
    )

    # Fallback: scanned PDF → OCR
    if is_scanned_pdf(documents):
        print(f"Scanned PDF detected: {filename} — running OCR pipeline")
        ocr_docs = ocr_scanned_pdf(file_path)
        if ocr_docs:
            for doc in ocr_docs:
                doc.metadata["is_engineering_drawing"] = is_drawing
            return ocr_docs
        # If OCR also fails, return whatever text was extracted
        print("OCR produced no output — using raw extracted text")

    for doc in documents:
        page_number = doc.metadata.get("page", 0)
        doc.metadata.update(
            build_metadata(file_path=file_path, file_type="pdf", page=page_number)
        )
        doc.metadata["is_engineering_drawing"] = is_drawing
        doc.metadata["page_count"] = len(documents)

    return documents


# =====================================================
# LOAD IMAGE USING VISUAL RAG
# =====================================================

def load_image(
    file_path
):

    extracted_text = extract_text_from_image(
        file_path
    )

    image_type = detect_document_type(
        file_path
    )

    vision_summary = analyze_image(
        file_path,
        """
Analyze this image.

Return:
1. Document Type
2. Key Information
3. Objects Detected
4. Tables (if any)
5. Charts (if any)
6. Engineering Drawing Details (if applicable)
7. Errors or Warnings (if visible)
"""
    )

    combined_content = f"""
IMAGE TYPE:
{image_type}

VISION SUMMARY:
{vision_summary}

OCR TEXT:
{extracted_text}
"""

    print("\n========== VISUAL RAG ==========")
    print(combined_content[:2000])
    print("================================\n")

    metadata = build_metadata(
        file_path=file_path,
        file_type="image"
    )

    metadata["image_type"] = image_type

    document = Document(
        page_content=combined_content,
        metadata=metadata
    )

    return [document]

# =====================================================
# LOAD DOCX
# =====================================================

def load_docx(file_path):
    doc = DocxDocument(file_path)

    # Extract paragraphs
    para_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    # Extract tables as markdown-style text
    table_sections = []
    for i, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            table_sections.append(f"[TABLE {i+1}]\n" + "\n".join(rows))

    combined = para_text
    if table_sections:
        combined += "\n\n" + "\n\n".join(table_sections)

    return [Document(
        page_content=combined,
        metadata=build_metadata(file_path, "docx"),
    )]

# =====================================================
# LOAD TXT
# =====================================================

def load_txt(
    file_path
):

    with open(
        file_path,
        "r",
        encoding="utf-8"
    ) as f:

        text = f.read()

    return [
        Document(
            page_content=text,
            metadata={
                "source": os.path.basename(
                    file_path
                ),
                "type": "txt"
            }
        )
    ]

# =====================================================
# LOAD CSV
# =====================================================

def load_csv(file_path):
    df = pd.read_csv(file_path)
    columns = ", ".join(df.columns.tolist())
    stats = df.describe(include='all').to_string()
    sample = df.head(20).to_markdown() if hasattr(df, 'to_markdown') else df.head(20).to_string()
    text = f"COLUMNS: {columns}\n\nDATA SAMPLE:\n{sample}\n\nSTATISTICS:\n{stats}"
    return [Document(page_content=text, metadata=build_metadata(file_path, "csv"))]


# =====================================================
# LOAD XLSX
# =====================================================

def load_xlsx(file_path):
    xl = pd.ExcelFile(file_path)
    documents = []
    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name)
        columns = ", ".join(df.columns.astype(str).tolist())
        sample = df.head(20).to_string()
        text = f"SHEET: {sheet_name}\nCOLUMNS: {columns}\n\n{sample}"
        meta = build_metadata(file_path, "xlsx")
        meta["sheet"] = sheet_name
        documents.append(Document(page_content=text, metadata=meta))
    return documents

# =====================================================
# LOAD PPTX
# =====================================================

def load_pptx(
    file_path
):

    prs = Presentation(
        file_path
    )

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(
                shape,
                "text"
            ):

                text += (
                    shape.text + "\n"
                )

    return [
        Document(
            page_content=text,
            metadata={
                "source": os.path.basename(
                    file_path
                ),
                "type": "pptx"
            }
        )
    ]

# =====================================================
# UNIVERSAL DOCUMENT LOADER
# =====================================================

def load_document(
    file_path
):

    extension = os.path.splitext(
        file_path
    )[1].lower()

    # -------------------------------------------------
    # PDF
    # -------------------------------------------------

    if extension == ".pdf":

        return load_pdf(
            file_path
        )
    if extension == ".docx":

        return load_docx(
            file_path
        )

    if extension == ".txt":

        return load_txt(
            file_path
        )

    if extension == ".csv":

        return load_csv(
            file_path
        )

    if extension == ".xlsx":

        return load_xlsx(
            file_path
        )

    if extension == ".pptx":

        return load_pptx(
            file_path
        )
    # -------------------------------------------------
    # IMAGES
    # -------------------------------------------------

    image_extensions = [
        ".png",
        ".jpg",
        ".jpeg"
    ]

    if extension in image_extensions:

        return load_image(
            file_path
        )

    # -------------------------------------------------
    # UNSUPPORTED
    # -------------------------------------------------

    raise ValueError(
        f"Unsupported file type: {extension}"
    )